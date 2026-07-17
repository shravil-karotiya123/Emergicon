import os
import math
import random
from PIL import Image, ImageDraw, ImageFilter

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def generate_background():
    """
    Generates a premium, high-definition 1920x1080 gradient mesh background
    with floating particles and neural network lines.
    """
    print("Generating premium background image...")
    
    # 1. Generate gradient mesh at a lower resolution to achieve smooth gradients
    # We will draw radial glows and then upscale with bilinear interpolation and blur.
    scale_factor = 10
    small_w = 1080 // scale_factor
    small_h = 1920 // scale_factor
    
    small_bg = Image.new("RGBA", (small_w, small_h), (3, 3, 3, 255))
    pixels = small_bg.load()
    
    # Centers of glows
    # Blue glow center (bottom-left area)
    bx, by = int(small_w * 0.25), int(small_h * 0.75)
    # Silver-white glow center (top-right area)
    wx, wy = int(small_w * 0.75), int(small_h * 0.25)
    
    for y in range(small_h):
        for x in range(small_w):
            # Blue glow calculation (#4FC3F7)
            dist_b = math.sqrt((x - bx)**2 + (y - by)**2)
            factor_b = max(0.0, 1.0 - dist_b / 120.0)  # Radius 120
            glow_b = int(factor_b * factor_b * 32)
            
            # White glow calculation
            dist_w = math.sqrt((x - wx)**2 + (y - wy)**2)
            factor_w = max(0.0, 1.0 - dist_w / 95.0)   # Radius 95
            glow_w = int(factor_w * factor_w * 22)
            
            # Blend base dark grey/black (3, 3, 3) with glows
            r = min(255, 3 + glow_w + int(glow_b * 0.31))
            g = min(255, 3 + glow_w + int(glow_b * 0.76))
            b = min(255, 3 + glow_w + glow_b)
            pixels[x, y] = (r, g, b, 255)
            
    # Resize to full-HD portrait 1080x1920 and apply heavy Gaussian Blur for smooth transition
    bg = small_bg.resize((1080, 1920), Image.Resampling.BILINEAR)
    bg = bg.filter(ImageFilter.GaussianBlur(25))
    
    # Create overlay drawing context for high-res details
    draw = ImageDraw.Draw(bg)
    
    # 2. Draw neural network/digital mesh lines (faint cyan lines connecting nodes)
    print("Drawing digital neural network...")
    random.seed(2026) # Match academic year for fun and determinism
    node_count = 16
    nodes = []
    
    # Generate random node positions
    for _ in range(node_count):
        nx = random.randint(100, 980)
        ny = random.randint(150, 1770)
        nodes.append((nx, ny))
        
    # Draw connections
    for i, node1 in enumerate(nodes):
        for node2 in nodes[i+1:]:
            dist = math.sqrt((node1[0] - node2[0])**2 + (node1[1] - node2[1])**2)
            if dist < 320:
                # Faint opacity based on distance
                opacity = int((1.0 - dist / 320.0) * 14)
                draw.line([node1, node2], fill=(79, 195, 247, opacity), width=1)
                
        # Draw a small subtle node ring
        draw.ellipse([node1[0]-3, node1[1]-3, node1[0]+3, node1[1]+3], outline=(255, 255, 255, 25), width=1)

    # 3. Draw floating particles (glowing white circles)
    print("Drawing floating particles...")
    for _ in range(75):
        px = random.randint(0, 1080)
        py = random.randint(0, 1920)
        radius = random.uniform(0.7, 2.8)
        opacity = random.randint(30, 160)
        
        # Soft outer glow ring
        draw.ellipse(
            [px - radius*2.2, py - radius*2.2, px + radius*2.2, py + radius*2.2], 
            fill=(255, 255, 255, int(opacity * 0.3))
        )
        # Sharp core
        draw.ellipse(
            [px - radius, py - radius, px + radius, py + radius], 
            fill=(255, 255, 255, opacity)
        )
        
    # Ensure assets directory exists
    os.makedirs("assets", exist_ok=True)
    bg_path = os.path.join("assets", "presentation_bg.png")
    bg.save(bg_path, "PNG")
    print(f"Background successfully saved to {bg_path}")
    return bg_path

def build_presentation(bg_path):
    """
    Assembles the PowerPoint presentation with widescreen dimensions,
    full-bleed background, centered clickable logo, and styled text.
    """
    print("Assembling PowerPoint presentation...")
    prs = Presentation()
    
    # 1. Set A4 Portrait dimensions (8.27 x 11.69 Inches)
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    
    # 2. Add a blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 3. Add full-bleed background picture
    slide.shapes.add_picture(bg_path, Inches(0), Inches(0), Inches(8.27), Inches(11.69))
    
    # 4. Center the EMERGICON logo
    logo_path = os.path.join("assets", "logo.png")
    if not os.path.exists(logo_path):
        raise FileNotFoundError(f"Logo not found at {logo_path}")
        
    # Load logo image to calculate aspect ratio
    with Image.open(logo_path) as logo_img:
        logo_w, logo_h = logo_img.size
    aspect = logo_h / logo_w
    
    # Define design sizes for portrait (make it slightly wider since portrait is narrower)
    target_w_in = 5.0
    target_h_in = target_w_in * aspect
    
    # Calculate centering coordinates (Shift slightly upwards to leave room for text)
    left_in = (8.27 - target_w_in) / 2
    top_in = 5.0 - (target_h_in / 2)
    
    print(f"Centering logo: left={left_in:.3f}\", top={top_in:.3f}\", width={target_w_in:.3f}\", height={target_h_in:.3f}\"")
    logo_shape = slide.shapes.add_picture(
        logo_path, 
        Inches(left_in), 
        Inches(top_in), 
        width=Inches(target_w_in), 
        height=Inches(target_h_in)
    )
    
    # 5. Add interactive hyperlink action to the logo shape
    # We will overlay a transparent rectangle directly on top of the logo shape.
    # This is the most reliable way to create a clickable area in pptx that supports URL hyperlinks.
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(target_w_in),
        Inches(target_h_in)
    )
    # Make overlay completely invisible (no fill, no line outline)
    overlay.fill.background()
    overlay.line.fill.background()
    
    # Apply Hyperlink click setting
    click_action = overlay.click_action
    click_action.hyperlink.address = "https://shravil-karotiya123.github.io/Emergicon/"
    print("Hyperlink target registered: https://shravil-karotiya123.github.io/Emergicon/")
    
    # 5.5 Add Header Text Above Logo
    header_top_in = top_in - 0.8
    header_box = slide.shapes.add_textbox(Inches(0), Inches(header_top_in), Inches(8.27), Inches(0.8))
    htf = header_box.text_frame
    hp = htf.paragraphs[0]
    hp.text = "Forum Installation Ceremony 2026-2027"
    hp.alignment = PP_ALIGN.CENTER
    hp.font.name = "Montserrat"
    hp.font.size = Pt(28)
    hp.font.color.rgb = RGBColor(255, 255, 255)
    hp.font.bold = True
    
    # 6. Add "Do not click on the logo" text box below
    text_top_in = top_in + target_h_in + 0.6
    text_left_in = 0
    text_w_in = 8.27
    text_h_in = 0.8
    
    tx_box = slide.shapes.add_textbox(
        Inches(text_left_in), 
        Inches(text_top_in), 
        Inches(text_w_in), 
        Inches(text_h_in)
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = "Don't\nclick on the logo"
    p.alignment = PP_ALIGN.CENTER
    
    # Set premium typography matching the brand identity
    p.font.name = "Montserrat"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(92, 188, 204) # #5CBCCC
    # Enable bold
    p.font.bold = True
    
    # Save the final .pptx presentation
    out_path = "emergicon_invitation_mobile.pptx"
    prs.save(out_path)
    print(f"Presentation successfully created and saved to {out_path}!")

if __name__ == "__main__":
    try:
        bg_path = generate_background()
        build_presentation(bg_path)
        print("PowerPoint generation workflow complete.")
    except Exception as e:
        print(f"An error occurred: {e}")
        import traceback
        traceback.print_exc()
